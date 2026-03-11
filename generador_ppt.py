import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ----------------------------
# PPT 생성
# - 제목: 항상 중앙 고정
# - 가사: KR 위 / ES 아래
# - UI 입력: KR 왼쪽 / ES 오른쪽
# ----------------------------
def crear_ppt(titulos_kr, bloques_kr, bloques_es, secuencia, estilos, resaltados):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, titulo in enumerate(titulos_kr):

        # ---------- 제목 슬라이드 (중앙 고정) ----------
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*estilos["bg_titulo"])

        tb = slide.shapes.add_textbox(
            Inches(1),
            Inches(3.2),              # ✅ 항상 중앙 근처 고정
            Inches(11.33),
            Inches(2),
        )
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = titulo
        r.font.size = Pt(estilos["tamano_titulo_kr"])
        r.font.color.rgb = RGBColor(*estilos["color_titulo_kr"])
        p.alignment = PP_ALIGN.CENTER

        # ---------- 가사 슬라이드 ----------
        for bloque_id in secuencia[i]:
            kr_lines = bloques_kr[i].get(bloque_id, [])
            es_lines = bloques_es[i].get(bloque_id, [])

            for j in range(len(kr_lines)):
                linea_kr = kr_lines[j]
                linea_es = es_lines[j] if j < len(es_lines) else ""

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(*estilos["bg_letra"])

                # KR (위)
                tb_kr = slide.shapes.add_textbox(
                    Inches(1),
                    Inches(estilos["altura_kr"]),
                    Inches(11.33),
                    Inches(1.6),
                )
                tf_kr = tb_kr.text_frame
                tf_kr.clear()
                tf_kr.word_wrap = True

                pkr = tf_kr.paragraphs[0]
                rkr = pkr.add_run()
                rkr.text = linea_kr
                rkr.font.size = Pt(estilos["tamano_letra_kr"])

                if bloque_id in resaltados[i]:
                    rkr.font.color.rgb = RGBColor(255, 192, 0)
                else:
                    rkr.font.color.rgb = RGBColor(*estilos["color_letra_kr"])

                pkr.alignment = PP_ALIGN.CENTER

                # ES (아래)
                if linea_es.strip():
                    tb_es = slide.shapes.add_textbox(
                        Inches(1),
                        Inches(estilos["altura_es"]),
                        Inches(11.33),
                        Inches(1.6),
                    )
                    tf_es = tb_es.text_frame
                    tf_es.clear()
                    tf_es.word_wrap = True

                    pes = tf_es.paragraphs[0]
                    res = pes.add_run()
                    res.text = linea_es
                    res.font.size = Pt(estilos["tamano_letra_es"])
                    res.font.color.rgb = RGBColor(*estilos["color_letra_es"])
                    pes.alignment = PP_ALIGN.CENTER

    return prs


# ----------------------------
# Streamlit UI
# ----------------------------
st.set_page_config(layout="wide")
st.title("피피티 잘 부탁드립니당~")

col1, col2, col3, col4 = st.columns(4)
with col1:
    num_canciones = st.number_input("찬양 개수", min_value=1, max_value=10, step=1)
with col2:
    size_titulo_kr = st.number_input("제목 글자 크기", value=36)
with col3:
    size_letra_kr = st.number_input("가사 한국어 글자 크기", value=36)
with col4:
    size_letra_es = st.number_input("가사 스페인어 글자 크기", value=28)

# ✅ 가사 위치만 조절 (제목 위치 슬라이더 제거됨)
pos1, pos2 = st.columns(2)
with pos1:
    altura_kr = st.slider("한국어 가사 위치 (PPT에서 위)", 0.0, 6.0, value=1.2, step=0.1)
with pos2:
    altura_es = st.slider("스페인어 가사 위치 (PPT에서 아래)", 0.0, 6.0, value=3.0, step=0.1)

# 색상 (고정)
color_titulo_kr = "#000000"
bg_titulo = "#FFFFFF"
color_letra_kr = "#FFFFFF"
color_letra_es = "#FFFFFF"
bg_letra = "#000000"

estilos = {
    "color_titulo_kr": tuple(int(color_titulo_kr[i:i+2], 16) for i in (1, 3, 5)),
    "bg_titulo": tuple(int(bg_titulo[i:i+2], 16) for i in (1, 3, 5)),
    "bg_letra": tuple(int(bg_letra[i:i+2], 16) for i in (1, 3, 5)),
    "color_letra_kr": tuple(int(color_letra_kr[i:i+2], 16) for i in (1, 3, 5)),
    "color_letra_es": tuple(int(color_letra_es[i:i+2], 16) for i in (1, 3, 5)),
    "altura_kr": altura_kr,
    "altura_es": altura_es,
    "tamano_titulo_kr": size_titulo_kr,
    "tamano_letra_kr": size_letra_kr,
    "tamano_letra_es": size_letra_es,
}

korean_titles = []
bloques_por_cancion_kr, bloques_por_cancion_es = [], []
secuencias, resaltados = [], []

for i in range(num_canciones):
    st.subheader(f"🎵 찬양 {i+1}")
    titulo = st.text_input(f"한국어 [제목] #{i+1}", key=f"kr_title_{i}")
    korean_titles.append(titulo)

    # UI는 좌/우 입력
    u1, u2 = st.columns(2)
    with u1:
        raw_kr = st.text_area("KR 전체 가사", key=f"kr_{i}", height=240)
    with u2:
        raw_es = st.text_area("ES 전체 가사", key=f"es_{i}", height=240)

    def parse(raw):
        blocks, cur = {}, None
        for line in raw.splitlines() + [""]:
            s = line.strip()
            if not s:
                cur = None
                continue
            if cur is None:
                cur = s
                blocks[cur] = []
            else:
                blocks[cur].append(s)
        return blocks

    bloques_kr = parse(raw_kr)
    bloques_es = parse(raw_es)

    bloques_por_cancion_kr.append(bloques_kr)
    bloques_por_cancion_es.append(bloques_es)

    secuencia_str = st.text_input("슬라이드 순서 (쉼표, 대문자/소문자, 띄어쓰기 없이 확인!)", key=f"seq_{i}")
    resaltado_str = st.text_input("후렴 블록 (쉼표, 대문자/소문자, 띄어쓰기 없이 확인!)", key=f"res_{i}")

    secuencias.append([s for s in secuencia_str.split(",") if s in bloques_kr])
    resaltados.append([s for s in resaltado_str.split(",") if s])

if st.button("완료!"):
    ppt = crear_ppt(
        korean_titles,
        bloques_por_cancion_kr,
        bloques_por_cancion_es,
        secuencias,
        estilos,
        resaltados,
    )

    ppt_path = "ppt_generado.pptx"
    ppt.save(ppt_path)

    with open(ppt_path, "rb") as f:
        st.download_button("PPT 다운로드", f, file_name=ppt_path)

    if os.path.exists(ppt_path):
        os.remove(ppt_path)



